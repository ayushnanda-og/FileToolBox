# pdf_tools.py
from pathlib import Path
from pypdf import PdfReader, PdfWriter

def merge_pdfs(paths, output):
    writer = PdfWriter()
    for path in paths:
        reader = PdfReader(str(path))
        for page in reader.pages:
            writer.add_page(page)
    with open(output, "wb") as f:
        writer.write(f)
    return output

def split_pdf(src, output_dir):
    reader = PdfReader(str(src))
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    outputs = []
    for i, page in enumerate(reader.pages, 1):
        writer = PdfWriter()
        writer.add_page(page)
        out_path = output_dir / f"{Path(src).stem}_page{i}.pdf"
        with open(out_path, "wb") as f:
            writer.write(f)
        outputs.append(out_path)
    return outputs

import tempfile, os, subprocess, shutil
from pathlib import Path

def compress_pdf(input_pdf, output_pdf, target_mb: float | None = None):
    """
    Compress PDF with Ghostscript.
    If target_mb is given, try qualities in order until size <= target_mb.
    """
    qualities = ["screen", "ebook", "printer", "prepress"]  # weakest → strongest
    if target_mb is None:
        qualities = ["ebook"]  # default single pass

    LOCAL_GS = Path(__file__).with_name("bin") / "gswin64c.exe"
    gs_path = str(LOCAL_GS) if LOCAL_GS.exists() else shutil.which("gswin64c") or shutil.which("gs")
    if not gs_path:
        raise RuntimeError("Ghostscript not found.")

    in_path = Path(input_pdf)
    out_path = Path(output_pdf)

    for q in qualities:
        # write to a temp file first
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmpf:
            tmp_name = tmpf.name

        cmd = [
            gs_path,
            "-sDEVICE=pdfwrite",
            "-dCompatibilityLevel=1.4",
            f"-dPDFSETTINGS=/{q}",
            "-dNOPAUSE",
            "-dQUIET",
            "-dBATCH",
            f"-sOutputFile={tmp_name}",
            str(in_path),
        ]
        subprocess.run(cmd, check=True)

        # If no target given, accept the first result (ebook)
        if target_mb is None:
            shutil.move(tmp_name, out_path)
            return out_path

        size_mb = os.path.getsize(tmp_name) / (1024 * 1024)
        if size_mb <= target_mb:
            shutil.move(tmp_name, out_path)
            return out_path
        else:
            os.remove(tmp_name)  # too large, try next quality

    # If we exit the loop, even strongest quality is still too big
    raise RuntimeError(
        "Could not reach target size. "
        "Try a larger target or optimize the PDF manually."
    )

from pdf2docx import Converter

def pdf_to_docx(input_pdf: str, output_docx: str):
    cv = Converter(input_pdf)
    cv.convert(output_docx, start=0, end=None)
    cv.close()

import pdfplumber
import pandas as pd

def pdf_to_excel(input_pdf: str, output_excel: str):
    with pdfplumber.open(input_pdf) as pdf:
        all_tables = []
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                all_tables.append(pd.DataFrame(table))

    if not all_tables:
        raise ValueError("No tables found in the PDF.")

    # Write all tables to one Excel file with multiple sheets
    with pd.ExcelWriter(output_excel, engine="openpyxl") as writer:
        for i, df in enumerate(all_tables):
            df.to_excel(writer, sheet_name=f"Table{i+1}", index=False, header=False)

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io

def pdf_to_pptx(input_pdf: str, output_pptx: str):
    pdf = fitz.open(input_pdf)
    prs = Presentation()

    # Set slide dimensions to match typical screen
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for page in pdf:
        pix = page.get_pixmap(dpi=150)
        img_bytes = pix.tobytes("png")
        image_stream = io.BytesIO(img_bytes)

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide
        slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(output_pptx)

from PIL import Image

def pdf_to_images(input_pdf: str, output_folder: str):
    pdf = fitz.open(input_pdf)
    output_paths = []

    for i, page in enumerate(pdf):
        pix = page.get_pixmap(dpi=150)
        output_path = os.path.join(output_folder, f"page_{i+1}.jpg")
        pix.save(output_path)
        output_paths.append(output_path)

    return output_paths


def images_to_pdf(image_paths: list[str], output_pdf: str):
    if not image_paths:
        raise ValueError("No images selected.")
    
    images = [Image.open(p).convert("RGB") for p in image_paths]
    images[0].save(output_pdf, save_all=True, append_images=images[1:])

from pypdf import PdfReader, PdfWriter
from pypdf.generic import RectangleObject

def rotate_pdf(input_pdf: str, output_pdf: str, angle: int = 90):
    reader = PdfReader(input_pdf)
    writer = PdfWriter()

    for page in reader.pages:
        page.rotate(angle)
        writer.add_page(page)

    with open(output_pdf, "wb") as f:
        writer.write(f)


def protect_pdf(input_pdf: str, output_pdf: str, password: str):
    reader = PdfReader(input_pdf)
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    writer.encrypt(password)
    with open(output_pdf, "wb") as f:
        writer.write(f)


def unlock_pdf(input_pdf: str, output_pdf: str, password: str):
    reader = PdfReader(input_pdf)
    if reader.is_encrypted:
        reader.decrypt(password)

    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    with open(output_pdf, "wb") as f:
        writer.write(f)

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from tempfile import NamedTemporaryFile

def add_text_watermark(input_pdf: str, output_pdf: str, watermark_text: str):
    reader = PdfReader(input_pdf)
    writer = PdfWriter()

    # Create temporary watermark PDF
    with NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
        c = canvas.Canvas(temp_file.name, pagesize=letter)
        c.setFont("Helvetica-Bold", 48)
        c.setFillAlpha(0.3)
        width, height = letter
        c.saveState()
        c.translate(width / 2, height / 2)
        c.rotate(45)
        c.drawCentredString(0, 0, watermark_text)
        c.restoreState()
        c.save()
        watermark_pdf = PdfReader(temp_file.name)
        watermark_page = watermark_pdf.pages[0]

    for page in reader.pages:
        page.merge_page(watermark_page)
        writer.add_page(page)

    with open(output_pdf, "wb") as f:
        writer.write(f)

import pytesseract
from pdf2image import convert_from_path
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

# Set the path to your Tesseract binary
pytesseract.pytesseract.tesseract_cmd = r"C:\Users\HP\FileToolbox\bin\Tesseract-OCR\tesseract.exe"

def ocr_pdf(input_pdf: str, output_pdf: str):
    images = convert_from_path(
        input_pdf,
        dpi=300,
        poppler_path=r"C:\Users\HP\FileToolbox\bin\poppler-24.08.0\Library\bin"
    )
    c = canvas.Canvas(output_pdf, pagesize=letter)

    for img in images:
        text = pytesseract.image_to_string(img)
        width, height = letter
        text_lines = text.split("\n")
        y = height - 50
        for line in text_lines:
            c.drawString(40, y, line)
            y -= 14
            if y < 40:
                c.showPage()
                y = height - 50
        c.showPage()

    c.save()
